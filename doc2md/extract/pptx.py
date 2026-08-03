"""PPTX extraction.

One Markdown section per slide, with shapes ordered top-to-bottom then
left-to-right so images land where they sit on the slide. Image locators are
slide numbers: `q3-deck-s03-img02`.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..config import Config
from ..model import (
    Block,
    Document,
    ImageNamer,
    ImageRef,
    slide_locator,
    slugify,
    title_from_path,
)

AUTONAME_RE = re.compile(r"^(picture|image|graphic|content placeholder)\s*\d*$", re.IGNORECASE)


def extract(path: Path, cfg: Config) -> Document:
    try:
        deck = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001 - python-pptx raises broadly
        raise RuntimeError(f"could not open PPTX: {exc}") from exc

    title = _deck_title(deck) or title_from_path(path)
    doc = Document(source=path, title=title, slug=slugify(title, path.stem or "document"))
    namer = ImageNamer(doc.slug)

    for number, slide in enumerate(deck.slides, 1):
        slide_title = _slide_title(slide)
        doc.add(Block(kind="heading", text=slide_title or f"Slide {number}", level=2))

        locator = slide_locator(number)
        for shape in _ordered_shapes(slide):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                doc.add(Block(kind="image", image=_picture(shape, namer.next(locator), doc)))
                continue
            if shape.has_table:
                rendered = _table_markdown(shape.table)
                if rendered:
                    doc.add(Block(kind="table", text=rendered))
                continue
            text = _shape_text(shape)
            # The title is already the heading; don't repeat it in the body.
            if text and text != slide_title:
                doc.add(Block(kind="text", text=text))

        notes = _notes(slide)
        if notes:
            doc.add(Block(kind="text", text=f"> **Speaker notes:** {notes}"))

    if not doc.blocks:
        doc.note("presentation contains no slides")
    return doc


def _deck_title(deck) -> str:
    try:
        return (deck.core_properties.title or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def _slide_title(slide) -> str:
    try:
        placeholder = slide.shapes.title
    except Exception:  # noqa: BLE001
        return ""
    if placeholder is None or not placeholder.has_text_frame:
        return ""
    return re.sub(r"\s+", " ", placeholder.text_frame.text).strip()


def _ordered_shapes(slide) -> list[Any]:
    """Shapes in reading order, tolerating shapes with no position."""

    def key(shape):
        top = shape.top if shape.top is not None else 0
        left = shape.left if shape.left is not None else 0
        return (top, left)

    return sorted(slide.shapes, key=key)


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    lines = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text:
            continue
        # Indented paragraphs are bullets in practice.
        prefix = "- " if paragraph.level or len(shape.text_frame.paragraphs) > 1 else ""
        lines.append(f"{prefix}{'  ' * paragraph.level}{text}" if prefix else text)
    return "\n".join(lines).strip()


def _picture(shape, image_id: str, doc: Document) -> ImageRef:
    caption = _picture_caption(shape)
    image = ImageRef(id=image_id, caption=caption)
    try:
        blob = shape.image.blob
        ext = shape.image.ext or "png"
    except Exception as exc:  # noqa: BLE001 - linked pictures have no blob
        image.reason = "linked (not embedded) image not available"
        doc.note(f"{image.id}: {exc}")
        return image

    image.ext = f".{str(ext).lstrip('.').lower()}"
    image.data = blob
    return image


def _picture_caption(shape) -> str:
    for value in (_alt_text(shape), getattr(shape, "name", "") or ""):
        value = (value or "").strip()
        if value and not AUTONAME_RE.match(value):
            return value
    return ""


def _alt_text(shape) -> str:
    """Alt text lives on the shape's XML description attribute."""
    try:
        return shape._element._nvXxPr.cNvPr.get("descr", "") or ""
    except Exception:  # noqa: BLE001
        return ""


def _notes(slide) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        text = slide.notes_slide.notes_text_frame.text
    except Exception:  # noqa: BLE001
        return ""
    return re.sub(r"\s+", " ", text).strip()


def _table_markdown(table) -> str:
    rows = []
    for row in table.rows:
        cells = [c.text.replace("|", "\\|").replace("\n", " ").strip() for c in row.cells]
        if any(cells):
            rows.append(cells)
    if not rows:
        return ""

    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = [c or f"col{i + 1}" for i, c in enumerate(rows[0])]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    lines += ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join(lines)
