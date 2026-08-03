"""PPTX extraction — one section per slide, shapes in reading order."""

from __future__ import annotations

import pytest

from doc2md.extract import extract


@pytest.fixture
def pptx_doc(docs, cfg):
    return extract(docs["pptx"], cfg)


def test_title_from_core_properties(pptx_doc):
    assert pptx_doc.title == "Launch Deck"


def test_each_slide_becomes_a_heading(pptx_doc):
    headings = [b.text for b in pptx_doc.blocks if b.kind == "heading"]
    assert headings == ["Agenda", "The Numbers"]


def test_slide_headings_are_level_two(pptx_doc):
    assert all(b.level == 2 for b in pptx_doc.blocks if b.kind == "heading")


def test_picture_is_extracted_with_slide_locator(pptx_doc):
    images = pptx_doc.images
    assert len(images) == 1
    assert images[0].extracted
    assert images[0].id.endswith("-s02-img01")  # second slide


def test_alt_text_becomes_the_caption(pptx_doc):
    assert pptx_doc.images[0].caption == "Sales photo"


def test_speaker_notes_are_included_and_marked(pptx_doc):
    notes = [b.text for b in pptx_doc.blocks if b.text.startswith("> **Speaker notes:**")]
    assert len(notes) == 1
    assert "Remember to pause here." in notes[0]


def test_title_text_is_not_repeated_in_the_body(pptx_doc):
    body = [b.text for b in pptx_doc.blocks if b.kind == "text"]
    assert "Agenda" not in body
    assert "The Numbers" not in body


def test_body_placeholder_text_is_kept(pptx_doc):
    body = " ".join(b.text for b in pptx_doc.blocks if b.kind == "text")
    assert "First item" in body


def test_corrupt_pptx_raises_runtime_error(isolated_cwd, cfg):
    broken = isolated_cwd / "broken.pptx"
    broken.write_bytes(b"not a zip")
    with pytest.raises(RuntimeError, match="could not open PPTX"):
        extract(broken, cfg)


def test_slide_without_a_title_gets_an_ordinal(isolated_cwd, cfg):
    from pptx import Presentation

    target = isolated_cwd / "untitled.pptx"
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])  # blank layout, no title
    deck.save(str(target))

    result = extract(target, cfg)
    headings = [b.text for b in result.blocks if b.kind == "heading"]
    assert headings == ["Slide 1"]


def test_table_on_a_slide_renders_as_markdown(isolated_cwd, cfg):
    from pptx import Presentation
    from pptx.util import Inches

    target = isolated_cwd / "tables.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Numbers"
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1))
    table = shape.table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Total"
    table.cell(1, 0).text = "North"
    table.cell(1, 1).text = "120"
    deck.save(str(target))

    doc = extract(target, cfg)
    rendered = [b.text for b in doc.blocks if b.kind == "table"]
    assert rendered
    assert "| Region | Total |" in rendered[0]
    assert "| North | 120 |" in rendered[0]


def test_multi_paragraph_placeholder_becomes_bullets(isolated_cwd, cfg):
    from pptx import Presentation

    target = isolated_cwd / "bullets.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Agenda"
    frame = slide.placeholders[1].text_frame
    frame.text = "First"
    frame.add_paragraph().text = "Second"
    deck.save(str(target))

    doc = extract(target, cfg)
    body = "\n".join(b.text for b in doc.blocks if b.kind == "text")
    assert "- First" in body
    assert "- Second" in body


def test_deck_without_core_title_falls_back_to_filename(isolated_cwd, cfg):
    from pptx import Presentation

    target = isolated_cwd / "untitled-deck.pptx"
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    deck.save(str(target))

    doc = extract(target, cfg)
    assert doc.title == "untitled deck"


def test_slide_without_notes_adds_nothing(isolated_cwd, cfg):
    from pptx import Presentation

    target = isolated_cwd / "nonotes.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Bare"
    deck.save(str(target))

    doc = extract(target, cfg)
    assert not any("Speaker notes" in b.text for b in doc.blocks)
