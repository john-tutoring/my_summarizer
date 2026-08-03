"""Programmatic builders for the sample documents the tests run against.

Nothing binary is committed. Each builder produces a file containing exactly
the structure its tests assert on, which keeps the assertions honest — when a
test says "the image sits between these two paragraphs", the fixture is the
reason that is true.

The ODT builder writes the zip by hand rather than shelling out to LibreOffice,
which does not exist on CI runners.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

# --- shared image helpers --------------------------------------------------


def png_bytes(color: tuple[int, int, int] = (30, 90, 200), size: tuple[int, int] = (200, 150)) -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", size, color).save(buf, "PNG")
    return buf.getvalue()


def jpg_bytes(color: tuple[int, int, int] = (200, 60, 40), size: tuple[int, int] = (200, 150)) -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", size, color).save(buf, "JPEG")
    return buf.getvalue()


# --- PDF -------------------------------------------------------------------


def make_pdf(path: Path) -> Path:
    """3 pages, 2 figures with captions, a repeated logo, a tiny icon, a TOC.

    Page 1: title (24pt), heading (16pt), body, figure 1, caption, body
    Page 2: heading, body before figure, figure 2, caption, body after
    Page 3: heading, body only
    """
    import pymupdf

    logo = png_bytes((200, 200, 200), (80, 30))
    icon = png_bytes((10, 10, 10), (8, 8))  # 64 px, below min_image_pixels
    fig1 = png_bytes((30, 90, 200))
    fig2 = png_bytes((200, 60, 40))

    doc = pymupdf.open()
    for n in range(3):
        page = doc.new_page()
        # On every page: caught by the furniture heuristic regardless of size.
        page.insert_image(pymupdf.Rect(450, 30, 530, 60), stream=logo)
        # On page 1 only, so the pixel floor is what excludes it — keeping the
        # two filters independently testable.
        if n == 0:
            page.insert_image(pymupdf.Rect(72, 40, 80, 48), stream=icon)

        if n == 0:
            page.insert_text((72, 90), "Annual Review", fontsize=24)
            page.insert_text((72, 130), "Introduction", fontsize=16)
            page.insert_text((72, 160), "Opening paragraph before the figure.", fontsize=11)
            page.insert_image(pymupdf.Rect(72, 180, 222, 280), stream=fig1)
            page.insert_text((72, 300), "Figure 1. Revenue by segment.", fontsize=9)
            page.insert_text((72, 330), "Closing paragraph after the figure.", fontsize=11)
        elif n == 1:
            page.insert_text((72, 90), "Results", fontsize=16)
            page.insert_text((72, 120), "Text before the second figure.", fontsize=11)
            page.insert_image(pymupdf.Rect(72, 140, 222, 240), stream=fig2)
            page.insert_text((72, 260), "Figure 2. Cost breakdown.", fontsize=9)
            page.insert_text((72, 290), "Text after the second figure.", fontsize=11)
        else:
            page.insert_text((72, 90), "Conclusion", fontsize=16)
            page.insert_text((72, 120), "Final remarks with no figures.", fontsize=11)

    doc.set_toc([[1, "Introduction", 1], [1, "Results", 2], [1, "Conclusion", 3]])
    doc.set_metadata({"title": "Annual Review 2026"})
    doc.save(str(path))
    doc.close()
    return path


def make_encrypted_pdf(path: Path) -> Path:
    import pymupdf

    doc = pymupdf.open()
    doc.new_page().insert_text((72, 90), "secret", fontsize=12)
    doc.save(str(path), encryption=pymupdf.PDF_ENCRYPT_AES_256, owner_pw="o", user_pw="u")
    doc.close()
    return path


# --- DOCX ------------------------------------------------------------------


def make_docx(path: Path, image_path: Path) -> Path:
    """Headings by style and by outlineLvl, an inline image, a caption, a table."""
    import docx
    from docx.oxml.ns import qn

    document = docx.Document()
    document.core_properties.title = "Quarterly Memo"

    document.add_heading("Overview", level=1)
    document.add_paragraph("Opening paragraph before the chart.")
    document.add_picture(str(image_path))
    caption = document.add_paragraph("Figure 1. Quarterly chart.")
    caption.style = document.styles["Caption"]
    document.add_paragraph("Paragraph after the chart.")

    document.add_heading("Detail", level=2)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Total"
    table.cell(1, 0).text = "North"
    table.cell(1, 1).text = "120"

    # A heading expressed only via outlineLvl, as localized Word builds do.
    para = document.add_paragraph("Outline Level Heading")
    p_pr = para._p.get_or_add_pPr()
    outline = p_pr.makeelement(qn("w:outlineLvl"), {qn("w:val"): "2"})
    p_pr.append(outline)

    document.add_paragraph("Closing text.")
    document.save(str(path))
    return path


# --- PPTX ------------------------------------------------------------------


def make_pptx(path: Path, image_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    deck.core_properties.title = "Launch Deck"

    first = deck.slides.add_slide(deck.slide_layouts[1])
    first.shapes.title.text = "Agenda"
    first.placeholders[1].text_frame.text = "First item"

    second = deck.slides.add_slide(deck.slide_layouts[5])
    second.shapes.title.text = "The Numbers"
    picture = second.shapes.add_picture(str(image_path), Inches(1), Inches(2), Inches(3), Inches(2))
    picture._element._nvXxPr.cNvPr.set("descr", "Sales photo")
    second.notes_slide.notes_text_frame.text = "Remember to pause here."

    deck.save(str(path))
    return path


# --- EPUB ------------------------------------------------------------------

_CH_WITH_HEADING = (
    "<html><body><h1>The Beginning</h1><p>Once upon a time.</p>"
    '<figure><img src="images/plate.png" alt="A green plate"/>'
    "<figcaption>Plate I</figcaption></figure><p>After the plate.</p></body></html>"
)
# No heading tag at all: forces the TOC-title lookup.
_CH_NO_HEADING = "<html><body><p>This chapter has no heading element of its own.</p></body></html>"
_CH_MISSING_IMAGE = (
    "<html><body><h1>The End</h1><p>Things conclude.</p>"
    '<img src="images/missing.png" alt="not here"/></body></html>'
)


def make_epub(path: Path) -> Path:
    """3 spine items; a nested TOC with real titles; one chapter lacking a heading."""
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier("id-test")
    book.set_title("A Short Book")
    book.set_language("en")

    book.add_item(
        epub.EpubItem(
            uid="im1",
            file_name="images/plate.png",
            media_type="image/png",
            content=png_bytes((90, 180, 90)),
        )
    )

    chapters = []
    for name, content in (
        ("c1.xhtml", _CH_WITH_HEADING),
        ("c2.xhtml", _CH_NO_HEADING),
        ("c3.xhtml", _CH_MISSING_IMAGE),
    ):
        item = epub.EpubHtml(title=name, file_name=name, content=content)
        book.add_item(item)
        chapters.append(item)

    # Nested: a Section wrapping the chapters, mirroring a real book's parts.
    book.toc = (
        epub.Link("c1.xhtml", "The Beginning", "ch1"),
        (
            epub.Section("Part One"),
            (
                epub.Link("c2.xhtml", "A Titled Chapter", "ch2"),
                epub.Link("c3.xhtml", "The End", "ch3"),
            ),
        ),
    )
    book.spine = chapters
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(path), book)
    return path


# --- OpenDocument ----------------------------------------------------------

_ODT_CONTENT = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content
  xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
  xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"
  xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
  xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0"
  xmlns:xlink="http://www.w3.org/1999/xlink">
 <office:body><office:text>
  <text:h text:outline-level="1">Overview</text:h>
  <text:p>Opening paragraph before the chart.</text:p>
  <text:p><draw:frame draw:name="Chart frame">
    <draw:image xlink:href="Pictures/chart.png"/></draw:frame></text:p>
  <text:p>Paragraph after the chart.</text:p>
  <text:h text:outline-level="2">Detail</text:h>
  <table:table>
   <table:table-row>
    <table:table-cell><text:p>Region</text:p></table:table-cell>
    <table:table-cell><text:p>Total</text:p></table:table-cell>
   </table:table-row>
   <table:table-row>
    <table:table-cell><text:p>North</text:p></table:table-cell>
    <table:table-cell><text:p>120</text:p></table:table-cell>
   </table:table-row>
  </table:table>
  <text:list><text:list-item><text:p>A list item</text:p></text:list-item></text:list>
  <text:p><draw:frame draw:name="Missing frame">
    <draw:image xlink:href="Pictures/gone.png"/></draw:frame></text:p>
 </office:text></office:body>
</office:document-content>
"""

_ODT_META = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-meta
  xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
  xmlns:dc="http://purl.org/dc/elements/1.1/">
 <office:meta><dc:title>Quarterly Memo</dc:title></office:meta>
</office:document-meta>
"""


def make_odt(path: Path) -> Path:
    """Hand-built ODT zip. Never shells out to LibreOffice."""
    with zipfile.ZipFile(path, "w") as archive:
        # mimetype must be first and stored uncompressed per the ODF spec.
        archive.writestr(
            zipfile.ZipInfo("mimetype"),
            "application/vnd.oasis.opendocument.text",
            compress_type=zipfile.ZIP_STORED,
        )
        archive.writestr("content.xml", _ODT_CONTENT)
        archive.writestr("meta.xml", _ODT_META)
        archive.writestr("Pictures/chart.png", png_bytes((40, 120, 200)))
    return path


# --- RTF -------------------------------------------------------------------


def make_rtf(path: Path) -> Path:
    """Control words, a font table and picture group to skip, a \\u escape."""
    content = (
        r"{\rtf1\ansi\deff0"
        r"{\fonttbl{\f0\froman Times New Roman;}{\f1\fswiss Arial;}}"
        r"{\colortbl;\red0\green0\blue0;}"
        r"\pard First paragraph of the document.\par "
        r"Second paragraph with a caf\u233\'3f accent.\par "
        r"{\pict\pngblip 89504e470d0a1a0a}"
        r"\pard Third paragraph after the picture.\par}"
    )
    path.write_text(content, encoding="latin-1")
    return path


# --- HTML / Markdown / text ------------------------------------------------


def make_html(path: Path, image_name: str = "chart.png") -> Path:
    path.write_text(
        "<html><head><title>The Web Page</title></head><body>\n"
        "<h1>Main Heading</h1>\n"
        "<p>Some <strong>bold</strong> text and a <a href='#x'>link</a>.</p>\n"
        f"<img src='{image_name}' alt='local chart'>\n"
        "<h2>Second Section</h2>\n"
        "<ul><li>one</li><li>two</li></ul>\n"
        "<img src='https://example.com/remote.png' alt='a remote image'>\n"
        "<img src='gone.png' alt='broken'>\n"
        "<img src='../outside.png' alt='traversal'>\n"
        "<script>alert('should be stripped')</script>\n"
        "</body></html>",
        encoding="utf-8",
    )
    return path


def make_markdown(path: Path, image_name: str = "fig1.png") -> Path:
    path.write_text(
        "# User Guide\n\n"
        "Intro paragraph.\n\n"
        f"![the first figure]({image_name})\n\n"
        "## Setup\n\n"
        "![a remote one](https://example.com/x.png)\n\n"
        "![missing file](nope.png)\n\n"
        "```python\n"
        "# this is a comment, not a heading\n"
        'print("hi")\n'
        "```\n",
        encoding="utf-8",
    )
    return path


def make_txt(path: Path) -> Path:
    path.write_text(
        "The Quick Report\n\nFirst paragraph here.\n\nSecond paragraph, longer than the first.\n",
        encoding="utf-8",
    )
    return path


def make_csv(path: Path) -> Path:
    path.write_text(
        'name,qty,note\nwidget,3,"has, comma"\ngadget,7,plain\n',
        encoding="utf-8",
    )
    return path


def make_xlsx(path: Path) -> Path:
    import xlsxwriter

    book = xlsxwriter.Workbook(str(path))
    sheet = book.add_worksheet("Data")
    for row, values in enumerate([["Region", "Total"], ["North", 120], ["South", 95]]):
        for col, value in enumerate(values):
            sheet.write(row, col, value)
    book.close()
    return path


# --- synthetic Markdown for docsum -----------------------------------------

_SENTENCE = (
    "Lorem ipsum dolor sit amet consectetur adipiscing elit sed do eiusmod tempor "
    "incididunt ut labore et dolore magna aliqua ut enim ad minim veniam quis. "
)
WORDS_PER_SENTENCE = 25
SENTENCES_PER_PARAGRAPH = 12
#: Exactly 300 words. Exported so tests can size a document against the
#: 20,000-word chapter threshold instead of guessing.
WORDS_PER_PARAGRAPH = WORDS_PER_SENTENCE * SENTENCES_PER_PARAGRAPH
PARAGRAPH = _SENTENCE * SENTENCES_PER_PARAGRAPH


def long_markdown(chapters: int = 5, paragraphs_each: int = 20, level: int = 1) -> str:
    """A multi-chapter document.

    Total words = chapters * paragraphs_each * WORDS_PER_PARAGRAPH. Callers that
    need the chaptered path must clear the threshold with that product.
    """
    marker = "#" * level
    names = ["Introduction", "Methods", "Results", "Discussion", "Conclusion",
             "Appendix", "Afterword", "Notes"]
    parts = []
    for index in range(chapters):
        name = names[index % len(names)]
        parts.append(f"{marker} {name}\n\n" + (PARAGRAPH + "\n\n") * paragraphs_each)
    return "".join(parts)


def flat_markdown(paragraphs: int = 200) -> str:
    """Long, but with no headings at all — drives the chunking path."""
    return (PARAGRAPH + "\n\n") * paragraphs
